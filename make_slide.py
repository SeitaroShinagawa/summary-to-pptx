
"""
executable on Windows (bug happens at text_frame.fit_text (editing font) on Linux)
usage: python make_slide.py --config example.config --img_dict images --out out.pptx
"""
from pptx import Presentation
from pptx.util import Cm, Pt
import configparser
from glob import glob
from argparse import ArgumentParser


def main(args):

    #load summary
    config = configparser.ConfigParser()
    config.read(args.config, encoding="utf-8")
   
    #load template
    prs = Presentation(args.template) 

    #slide settings
    slide_layout = prs.slide_layouts[1] #bullet type(1 title, 1 list)
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes

    shapes = prs.slides[0].shapes #use slide number 0 (first slide)

    title = shapes[0] #title text box
    main = shapes[1] #list text box

    #edit title text box
    tf = title.text_frame
    tf.text = config["info"]["title"]
    tf.fit_text(font_family='Segoe UI', max_size=28, bold=True, italic=False, font_file=None)

    p = tf.add_paragraph()
    p.text = config["info"]["authors"]
    p.font.size = Pt(18)

    p = tf.add_paragraph()
    p.text = config["info"]["affiliation"]
    p.font.size = Pt(18)

    #edit list text box
    tf = main.text_frame
    for i,key in enumerate(sorted(config["main"].keys())):
        if i==0:
            tf.text = config["main"][key]
            tf.paragraphs[0].font.size = Pt(24)

        else:
            p = tf.add_paragraph()
            p.text = config["main"][key]
            p.font.size = Pt(24)

    #width = prs.slide_width
    #height = prs.slide_height

    pic_left = Cm(1.25)
    pic_top = Cm(12.9)
    pic_width = Cm(5.80)

    #image paste
    img_pathes = glob(args.img_dict+"/*.png")
    num = len(img_pathes)
    accum_left = 0
    for i,img_path in enumerate(img_pathes):
        pic = shapes.add_picture(img_path, pic_left+accum_left, pic_top, width=pic_width) 
        #pic.left = int( ( width - pic.width ) /num )
        #pic.top = int( ( height - pic.height ) /num )
        accum_left += pic.width

    prs.save(args.out)


if __name__ == "__main__":
    
    parser = ArgumentParser() 
    parser.add_argument('--config', '-c', default="example.config", help="config path (paper summary)")
    parser.add_argument('--template', '-t', default="template.pptx", help="template slide")
    parser.add_argument('--img_dict', '-i', default="images", help='output slide') 
    parser.add_argument('--out', '-o', default="out.pptx", help='output slide') 
    args = parser.parse_args()
    print(args)

    main(args)

